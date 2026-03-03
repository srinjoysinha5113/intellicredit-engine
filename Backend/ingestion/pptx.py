"""
PowerPoint file ingestion.
Handles .pptx format extraction.
"""

from typing import List
from pptx import Presentation

from config import CHUNK_SIZE, CHUNK_OVERLAP
from ingestion.chunking import split_text_into_chunks
from models.schemas import KnowledgeChunk


def process_pptx_file(file_path: str, file: str, relative_path: str, category: str, subcategory: str) -> List[KnowledgeChunk]:
    """Process PowerPoint file and create KnowledgeChunk objects."""
    chunks = []
    
    try:
        presentation = Presentation(file_path)
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = []
            slide_title = ""
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Check if it's likely a title (usually first or largest text)
                    if not slide_title and len(shape.text.strip()) < 100:
                        slide_title = shape.text.strip()
                    slide_text.append(shape.text.strip())
            
            # Combine slide text
            if slide_text:
                full_slide_text = f"Slide {slide_num + 1}"
                if slide_title:
                    full_slide_text += f": {slide_title}"
                full_slide_text += "\n" + "\n".join(slide_text)
                
                # Split into chunks
                slide_chunks = split_text_into_chunks(full_slide_text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP)
                
                for chunk in slide_chunks:
                    knowledge_chunk = KnowledgeChunk(
                        content=chunk,
                        source_type="pptx",
                        source_name=file,
                        source_path=relative_path,
                        metadata={"category": category, "subcategory": subcategory, "slide": slide_num + 1, "title": slide_title},
                        structure={"slide": slide_num + 1, "title": slide_title}
                    )
                    chunks.append(knowledge_chunk)
        
    except Exception as e:
        print(f"Error processing PowerPoint file {file}: {str(e)}")
    
    return chunks
